"""
Generate presentation slides for Robust Offline RL Disentanglement project.
Run: python docs/make_slides.py
Output: docs/presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os
import io
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches

# ---------------------------------------------------------------------------
# Color palette (clean academic style)
# ---------------------------------------------------------------------------
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy
C_BLUE   = RGBColor(0x16, 0x21, 0x3E)   # slide body bg
C_ACCENT = RGBColor(0x0F, 0x9B, 0xCF)   # bright blue accent
C_GOLD   = RGBColor(0xF5, 0xA6, 0x23)   # gold highlight
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xE8, 0xF4, 0xFD)   # very light blue text area
C_GREY   = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

FIGURES_ROOT = os.path.join(os.path.dirname(__file__), "..", "results", "figures")
NB_FIGURES   = os.path.join(os.path.dirname(__file__), "nb_figures")

def fig(rel_path):
    p = os.path.join(FIGURES_ROOT, rel_path)
    return p if os.path.exists(p) else None

def nbfig(filename):
    p = os.path.join(NB_FIGURES, filename)
    return p if os.path.exists(p) else None

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, fill_color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_image(slide, path, l, t, w, h=None):
    if path is None or not os.path.exists(path):
        return
    if h is None:
        slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    else:
        slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def slide_header(slide, title, subtitle=None):
    """Dark accent bar at top with title."""
    add_rect(slide, 0, 0, 13.33, 1.2, C_ACCENT)
    add_text(slide, title, 0.4, 0.15, 12.5, 0.9,
             font_size=28, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.85, 12.5, 0.4,
                 font_size=14, color=C_LIGHT, align=PP_ALIGN.LEFT)

def render_formula(latex_str, fontsize=18, bg="#1A1A2E", fg="white",
                   fig_w=8.0, fig_h=0.9):
    """
    Render a LaTeX-style formula string via matplotlib mathtext.
    Returns a BytesIO PNG stream suitable for add_picture.
    """
    fig, ax = plt.subplots(figsize=(fig_w, fig_h))
    fig.patch.set_facecolor(bg)
    ax.set_facecolor(bg)
    ax.axis("off")
    ax.text(0.0, 0.5, latex_str, fontsize=fontsize, color=fg,
            va="center", ha="left", transform=ax.transAxes,
            fontfamily="DejaVu Sans")
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                facecolor=fig.get_facecolor(), dpi=150)
    plt.close(fig)
    buf.seek(0)
    return buf


def add_formula(slide, latex_str, l, t, w, fontsize=17, bg="#1A1A2E", fg="white", fig_h=0.85):
    """Render formula and place it on the slide at position (l, t) with width w inches."""
    buf = render_formula(latex_str, fontsize=fontsize, bg=bg, fg=fg,
                         fig_w=w * 1.5, fig_h=fig_h)
    slide.shapes.add_picture(buf, Inches(l), Inches(t), width=Inches(w))


def bullet_box(slide, items, l, t, w, h,
               font_size=16, color=C_WHITE, title=None, title_color=C_GOLD):
    """Render a list of bullet strings into a single textbox."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = title
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = title_color
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color

# ---------------------------------------------------------------------------
# Build presentation
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank

# ===========================================================================
# Slide 1 — Title
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)

# decorative accent bar left
add_rect(slide, 0, 0, 0.18, 7.5, C_ACCENT)
# decorative thin line
add_rect(slide, 0.18, 3.0, 13.15, 0.06, C_GOLD)

add_text(slide, "Robust Offline RL", 0.7, 1.6, 12.0, 1.1,
         font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "via Disentangled Privileged Pretraining", 0.7, 2.6, 12.0, 0.9,
         font_size=30, bold=False, color=C_ACCENT, align=PP_ALIGN.LEFT)
add_text(slide, "Representation robustness under synthetic observation corruption\n"
                "Privileged Pretraining Framework (PPF) + disentanglement regularizers",
         0.7, 3.3, 11.5, 1.2,
         font_size=17, color=C_GREY, align=PP_ALIGN.LEFT)
add_text(slide, "Course Final Presentation  ·  2026", 0.7, 6.6, 10.0, 0.6,
         font_size=14, color=C_GREY, align=PP_ALIGN.LEFT)

# ===========================================================================
# Slide 2 — Motivation
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Motivation", "Why does robustness matter in offline RL?")

# Left column — problem
add_rect(slide, 0.4, 1.35, 5.8, 5.7, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["Real-world sensors are noisy or unreliable",
     "Offline datasets are collected under ideal conditions",
     "At deployment, observations can be corrupted",
     "Standard offline RL policies degrade sharply under noise"],
    0.6, 1.55, 5.4, 3.5,
    font_size=16, title="The Problem")

# Right column — challenge
add_rect(slide, 6.6, 1.35, 6.3, 5.7, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["Cannot collect more data online to adapt",
     "Encoder must separate task-relevant from noise features",
     "Three corruption families studied:\n"
     "   concat  |  project  |  nonlinear"],
    6.8, 1.55, 5.9, 3.5,
    font_size=16, title="The Challenge")

add_text(slide,
         "Goal: learn a robust state encoder from offline data that transfers to corrupted observations at test time",
         0.4, 6.2, 12.5, 0.9,
         font_size=16, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

# ===========================================================================
# Slide 3 — Problem Setup
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Problem Setup")

# Three noise type boxes
col_w, col_h = 3.6, 3.2
tops = 1.5
for i, (label, desc, color) in enumerate([
    ("concat",    "Append nuisance noise\ndirectly to clean obs\n\n→ dimension increases\n→ easy to detect", RGBColor(0x0F, 0x6E, 0x9B)),
    ("project",   "Random orthogonal linear\nmixing after concatenation\n\n→ features entangled\n→ moderate difficulty",    RGBColor(0x09, 0x55, 0x7A)),
    ("nonlinear", "Two-layer nonlinear\nmixing after concatenation\n\n→ highly entangled\n→ hardest setting",   RGBColor(0x06, 0x3A, 0x57)),
]):
    lft = 0.4 + i * (col_w + 0.35)
    add_rect(slide, lft, tops, col_w, col_h, color)
    add_text(slide, label, lft + 0.15, tops + 0.1, col_w - 0.3, 0.5,
             font_size=20, bold=True, color=C_GOLD)
    add_text(slide, desc, lft + 0.15, tops + 0.65, col_w - 0.3, col_h - 0.8,
             font_size=15, color=C_WHITE)

# Environments row
add_rect(slide, 0.4, 5.0, 12.5, 1.8, RGBColor(0x0D, 0x1B, 0x33))
add_text(slide, "Environments (D4RL medium-v2)", 0.6, 5.1, 6.0, 0.4,
         font_size=16, bold=True, color=C_GOLD)
envs = ["halfcheetah-medium-v2", "hopper-medium-v2", "walker2d-medium-v2", "ant-medium-v2"]
env_text = "   |   ".join(envs)
add_text(slide, env_text, 0.6, 5.55, 12.0, 0.5,
         font_size=15, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, "Policy: IQL (main)   +   TD3+BC / BC (ablation A)", 0.6, 6.1, 12.0, 0.5,
         font_size=14, color=C_GREY, align=PP_ALIGN.CENTER)

# ===========================================================================
# Slide 4 — PPF Method  (full-slide diagram)
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Privileged Pretraining Framework (PPF)",
             "Encoder trained with clean-state supervision; frozen at inference on noisy observations only")

# Embed the PPF architecture diagram (full width)
ppf_diagram = os.path.join(os.path.dirname(__file__), "ppf_diagram.png")
if os.path.exists(ppf_diagram):
    slide.shapes.add_picture(ppf_diagram,
                             Inches(0.25), Inches(1.3),
                             width=Inches(12.85))

# Regularizer strip at the bottom
add_rect(slide, 0.25, 6.85, 12.85, 0.55, RGBColor(0x08, 0x2C, 0x4A))
add_text(slide,
         "Regularizers:  Barlow Twins  |  Cov-Whitening  |  HSIC  |  Distance Corr  |  InfoNCE  |  L1",
         0.4, 6.9, 12.5, 0.42,
         font_size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

# ===========================================================================
# Slide 5 — Baselines
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Methods & Baselines")

rows = [
    ("true_only",              "Upper bound",      "Policy trained directly on clean states"),
    ("raw_noisy",              "Lower bound",      "Policy trained on raw corrupted observations"),
    ("plain (PPF)",            "PPF no reg.",      "PPF encoder — dynamics + reward, no disentanglement"),
    ("disentangled_barlow",    "PPF method",       "PPF + Barlow Twins cross-correlation penalty"),
    ("disentangled_cov",       "PPF method",       "PPF + covariance whitening penalty"),
    ("disentangled_hsic",      "PPF method",       "PPF + HSIC independence criterion"),
    ("disentangled_dcor",      "PPF method",       "PPF + distance correlation penalty"),
    ("disentangled_infonce",   "PPF method",       "PPF + InfoNCE contrastive penalty"),
    ("disentangled_l1",        "PPF method",       "PPF + L1 cross-correlation penalty"),
    ("PCA-IQL",                "External baseline","PCA projection (no neural encoder, no privileged info)"),
]

row_colors = {
    "Upper bound":      RGBColor(0x0A, 0x5C, 0x2E),
    "Lower bound":      RGBColor(0x6B, 0x1A, 0x1A),
    "PPF no reg.":      RGBColor(0x12, 0x3A, 0x5E),
    "PPF method":       RGBColor(0x0D, 0x2C, 0x4F),
    "External baseline":RGBColor(0x3A, 0x2A, 0x00),
}

col_x = [0.4, 3.5, 5.5]
col_w_vals = [3.0, 2.0, 7.2]
header_y = 1.4
add_rect(slide, 0.4, header_y, 12.5, 0.38, C_ACCENT)
for ci, label in enumerate(["Method", "Type", "Description"]):
    add_text(slide, label, col_x[ci] + 0.1, header_y + 0.04, col_w_vals[ci], 0.3,
             font_size=13, bold=True, color=C_WHITE)

for ri, (method, mtype, desc) in enumerate(rows):
    y = header_y + 0.38 + ri * 0.46
    add_rect(slide, 0.4, y, 12.5, 0.44, row_colors[mtype])
    for ci, text in enumerate([method, mtype, desc]):
        add_text(slide, text, col_x[ci] + 0.1, y + 0.06, col_w_vals[ci], 0.35,
                 font_size=12, color=C_WHITE)

# ===========================================================================
# Slide 6 — Robustness: Degradation Curve  (notebook cell 9)
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Results: Performance Degradation Curve",
             "hopper-medium-v2 · nonlinear noise · score vs. noise scale · IQL")

add_image(slide, nbfig("cell09_out01.png"), 0.3, 1.35, 12.7, 4.5)

add_rect(slide, 0.3, 6.0, 12.7, 1.1, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["dCor and Barlow show the flattest curves → most graceful degradation under increasing noise",
     "plain PPF already outperforms raw_noisy; disentanglement regularizers provide an additional margin at high scale",
     "PCA-IQL degrades sharply beyond scale 1.0"],
    0.5, 6.05, 12.3, 1.0,
    font_size=14)

# ===========================================================================
# Slide 7 — Score Heatmap   (notebook cell 13)
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Results: Full Score Heatmap",
             "hopper-medium-v2 · nonlinear · normalized score across all (dim, scale) configs")

add_image(slide, nbfig("cell13_out01.png"), 0.2, 1.35, 12.9, 4.6)

add_rect(slide, 0.2, 6.1, 12.9, 1.0, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["Top-left cells (small dim/scale) are easy — all methods score well",
     "Bottom-right corner (large dim & scale) is the hardest regime — disentangled methods retain color while plain/raw_noisy fade",
     "Heatmap reveals method-specific strengths: dCor / Barlow maintain brightness across the entire grid"],
    0.4, 6.15, 12.5, 0.9,
    font_size=14)

# ===========================================================================
# Slide 8 — Composite Ranking + Radar  (notebook cells 16 & 17)
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Method Selection: Composite Ranking",
             "hopper-medium-v2 · four metrics weighted: global mean 30%, hard-condition 35%, drop% 20%, win-rate 15%")

# Left: composite bar chart
add_image(slide, nbfig("cell16_out00.png"), 0.3, 1.4, 7.5, 3.6)
add_text(slide, "Composite score (higher = better overall)",
         0.3, 5.0, 7.5, 0.35, font_size=12, color=C_GREY, align=PP_ALIGN.CENTER)

# Right: radar chart
add_image(slide, nbfig("cell17_out00.png"), 8.0, 1.4, 5.0, 3.6)
add_text(slide, "Radar: top-4 methods vs plain — four evaluation axes",
         8.0, 5.0, 5.0, 0.35, font_size=12, color=C_GREY, align=PP_ALIGN.CENTER)

add_rect(slide, 0.3, 5.5, 12.6, 1.6, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["dCor ranks #1 overall: leads on both global mean and hard-condition robustness",
     "Barlow Twins #2: comparable hard-condition performance, slightly lower global mean",
     "plain PPF #3 — disentanglement consistently provides a meaningful uplift",
     "Verdict: dCor and Barlow are the recommended methods for high-noise regimes"],
    0.5, 5.6, 12.2, 1.4,
    font_size=14)

# ===========================================================================
# Slide 8b — Win Rate  (notebook cell 11)
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Results: Win Rate vs. Plain Baseline",
             "hopper-medium-v2 · nonlinear · fraction of (dim, scale) configs where method beats plain")

add_image(slide, nbfig("cell11_out03.png"), 0.3, 1.4, 12.7, 4.4)

add_rect(slide, 0.3, 6.0, 12.7, 1.1, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["dCor wins in 71.5% of configurations with an average margin of +2.8 pts",
     "Barlow wins in 69.4% — both methods reliably outperform the un-regularized PPF encoder",
     "InfoNCE and L1 show lower win rates, suggesting sensitivity to specific noise regimes"],
    0.5, 6.05, 12.3, 1.0,
    font_size=14)

# ===========================================================================
# Slide 9 — Ablation
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Ablation Study",
             "B1: remove privileged target   |   B2: reward-only pretraining")

# Two ablation figures side by side
fig_b1 = fig("main_ablation/halfcheetah-medium-v2/nonlinear/bar/scale_sweep_nonlinear_dim_17.png")
add_image(slide, fig_b1, 0.4, 1.45, 6.2, 3.5)
add_text(slide, "Full PPF vs no-privilege (B1) — halfcheetah nonlinear",
         0.4, 4.95, 6.2, 0.4, font_size=12, color=C_GREY, align=PP_ALIGN.CENTER)

fig_b2 = fig("main_ablation/halfcheetah-medium-v2/concat/bar/scale_sweep_concat_dim_40.png")
add_image(slide, fig_b2, 6.9, 1.45, 6.0, 3.5)
add_text(slide, "Full PPF vs reward-only (B2) — halfcheetah concat",
         6.9, 4.95, 6.0, 0.4, font_size=12, color=C_GREY, align=PP_ALIGN.CENTER)

# Findings row
add_rect(slide, 0.4, 5.5, 12.5, 1.65, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["B1 (no privileged supervision):  removing clean-state targets degrades all methods — confirms that privileged info is the key ingredient",
     "B2 (reward only):  dropping dynamics loss hurts further — next-state prediction provides a richer self-supervised signal than reward alone"],
    0.6, 5.6, 12.0, 1.5,
    font_size=14)

# ===========================================================================
# Slide — Noise Construction Formulas
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Noise Construction: Three Corruption Families",
             "How synthetic observation noise is injected into clean states")

# Row layout: three panels
panel_w, panel_h = 3.9, 4.9
panel_y = 1.4
for pi, (noise_type, color_rgb) in enumerate([
    ("concat",    RGBColor(0x0F, 0x6E, 0x9B)),
    ("project",   RGBColor(0x09, 0x55, 0x7A)),
    ("nonlinear", RGBColor(0x06, 0x3A, 0x57)),
]):
    lx = 0.35 + pi * (panel_w + 0.2)
    add_rect(slide, lx, panel_y, panel_w, panel_h, color_rgb)
    add_text(slide, noise_type, lx + 0.15, panel_y + 0.1, panel_w - 0.3, 0.45,
             font_size=20, bold=True, color=C_GOLD)

# --- concat formulas ---
lx0 = 0.35
add_formula(slide,
    r"$\varepsilon \sim \mathcal{N}(0,\,\sigma^2 I_d)$",
    lx0 + 0.1, panel_y + 0.65, panel_w - 0.25, fontsize=16, fig_h=0.7,
    bg="#0F6E9B")
add_formula(slide,
    r"$\tilde{o} = [s_{norm} \| \varepsilon]$",
    lx0 + 0.1, panel_y + 1.35, panel_w - 0.25, fontsize=16, fig_h=0.7,
    bg="#0F6E9B")
add_text(slide,
    "Clean state s is normalized;\nGaussian noise appended directly.\n"
    "Total dim = obs_dim + noise_dim.\n\n"
    "Corrupted dims are linearly\nseparable in principle.",
    lx0 + 0.12, panel_y + 2.1, panel_w - 0.25, 2.6,
    font_size=13, color=C_LIGHT)

# --- project formulas ---
lx1 = 0.35 + panel_w + 0.2
add_formula(slide,
    r"$\tilde{o}_{pre} = [s_{norm} \| \varepsilon]$",
    lx1 + 0.1, panel_y + 0.65, panel_w - 0.25, fontsize=16, fig_h=0.7,
    bg="#09557A")
add_formula(slide,
    r"$\tilde{o} = \tilde{o}_{pre} \cdot Q$",
    lx1 + 0.1, panel_y + 1.35, panel_w - 0.25, fontsize=16, fig_h=0.7,
    bg="#09557A")
add_text(slide,
    "Q is a fixed random orthogonal\nmatrix (QR decomp of Gaussian).\n\n"
    "State and noise are linearly\nmixed — requires linear algebra\nto separate.",
    lx1 + 0.12, panel_y + 2.1, panel_w - 0.25, 2.6,
    font_size=13, color=C_LIGHT)

# --- nonlinear formulas ---
lx2 = 0.35 + 2 * (panel_w + 0.2)
add_formula(slide,
    r"$h = \tanh(\tilde{o}_{pre} \cdot W_1)$",
    lx2 + 0.1, panel_y + 0.65, panel_w - 0.25, fontsize=15, fig_h=0.7,
    bg="#063A57")
add_formula(slide,
    r"$\tilde{o} = h \cdot W_2$",
    lx2 + 0.1, panel_y + 1.35, panel_w - 0.25, fontsize=16, fig_h=0.7,
    bg="#063A57")
add_text(slide,
    "W₁, W₂ are fixed random\northogonal matrices.\n\n"
    "tanh introduces nonlinearity —\nfeatures are deeply entangled.\n"
    "Hardest corruption to undo.",
    lx2 + 0.12, panel_y + 2.1, panel_w - 0.25, 2.6,
    font_size=13, color=C_LIGHT)

# Bottom note
add_rect(slide, 0.35, 6.45, 12.6, 0.7, RGBColor(0x0D, 0x1B, 0x33))
add_text(slide,
    "noise_dim d ∈ {4, 8, 13, 17}   |   noise_scale σ ∈ {0.5, 1.0, 1.5, 2.0}   |   "
    "pure_obs = s_norm (first obs_dim dims) — used as privileged supervision target",
    0.5, 6.52, 12.2, 0.55, font_size=13, color=C_GREY, align=PP_ALIGN.CENTER)

# ===========================================================================
# Slide — Disentanglement Regularizer Formulas
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Disentanglement Regularizers: Mathematical Formulas",
             "All methods minimize independence between z_task and z_irrel via different criteria")

# Total loss banner
add_rect(slide, 0.35, 1.38, 12.6, 0.72, RGBColor(0x08, 0x2C, 0x4A))
add_formula(slide,
    r"$\mathcal{L}_{total} = \mathcal{L}_{dyn} + \mathcal{L}_{rew} + \lambda \cdot \mathcal{L}_{reg}(z_{task},\; z_{irrel})$",
    0.5, 1.42, 12.1, fontsize=18, fig_h=0.62, bg="#082C4A")

# Six method boxes in 2×3 grid
methods = [
    ("Barlow Twins",
     r"$\tilde{z}_i = (z_i - \mu_i)/\sigma_i,\quad C = \tilde{z}_1^T \tilde{z}_2 / N$",
     r"$\mathcal{L}_{barlow} = \sum_{i,j} C_{ij}^2$",
     "Drives entire cross-corr matrix -> 0\n(L2 penalty, sensitive to large off-diag)"),
    ("Cov-Whitening",
     r"$\tilde{z}_i = (z_i - \mu_i)/\sigma_i,\quad C = \tilde{z}_1^T \tilde{z}_2 / N$",
     r"$\mathcal{L}_{cov} = \sum_{i,j} C_{ij}^2$",
     "Same matrix form as Barlow;\ndiffers in normalization + context"),
    ("HSIC",
     r"$K_{ij}=\exp(-\|z_1^i - z_1^j\|^2 / 2\sigma^2)$",
     r"$\mathcal{L}_{hsic} = \mathrm{tr}(K_c L_c) / (N-1)^2$",
     "Kernel-based independence test;\nK_c, L_c are double-centered kernels"),
    ("Distance Corr (dCor)",
     r"$A = \mathrm{dcenter}(\|z_1^i - z_1^j\|_2),\quad B = \mathrm{dcenter}(\|z_2^i - z_2^j\|_2)$",
     r"$\mathcal{L}_{dcor} = \overline{A \odot B} / \sqrt{\overline{A^2} \cdot \overline{B^2}}$",
     "Captures nonlinear dependence;\nzero iff z_task is independent of z_irrel"),
    ("InfoNCE Repulsion",
     r"$\hat{z}_i = z_i / \|z_i\|_2$",
     r"$\mathcal{L}_{rep} = \mathrm{mean}(\exp(\hat{z}_1 \hat{z}_2^T / \tau)),\quad \tau=0.5$",
     "Pushes representations apart\nvia cosine-similarity repulsion"),
    ("L1 Cross-Corr",
     r"$\tilde{z}_i = (z_i - \mu_i)/\sigma_i,\quad C = \tilde{z}_1^T \tilde{z}_2 / N$",
     r"$\mathcal{L}_{L1} = \sum_{i,j} |C_{ij}|$",
     "L1 penalty encourages sparse\ncorrelations (more robust to outliers)"),
]

cols, rows_n = 3, 2
box_w = 4.05
box_h = 2.28
start_x, start_y = 0.35, 2.22
gap_x, gap_y = 0.12, 0.1

for idx, (title, formula1, formula2, note) in enumerate(methods):
    col = idx % cols
    row = idx // cols
    lx = start_x + col * (box_w + gap_x)
    ly = start_y + row * (box_h + gap_y)
    add_rect(slide, lx, ly, box_w, box_h, RGBColor(0x0D, 0x1B, 0x33))
    add_text(slide, title, lx + 0.12, ly + 0.06, box_w - 0.2, 0.35,
             font_size=14, bold=True, color=C_GOLD)
    add_formula(slide, formula1, lx + 0.08, ly + 0.42, box_w - 0.15,
                fontsize=12, fig_h=0.52, bg="#0D1B33")
    add_formula(slide, formula2, lx + 0.08, ly + 0.96, box_w - 0.15,
                fontsize=14, fig_h=0.55, bg="#0D1B33")
    add_text(slide, note, lx + 0.12, ly + 1.58, box_w - 0.2, 0.65,
             font_size=11, color=C_GREY, italic=True)

# ===========================================================================
# Slide 9 — Conclusion
# ===========================================================================
slide = prs.slides.add_slide(blank_layout)
set_bg(slide, C_DARK)
slide_header(slide, "Conclusion & Future Work")

add_rect(slide, 0.4, 1.35, 12.5, 3.3, RGBColor(0x0D, 0x1B, 0x33))
bullet_box(slide,
    ["PPF (Privileged Pretraining Framework) effectively protects offline RL policies from observation corruption",
     "Disentanglement regularizers (Barlow Twins, HSIC, dCor, InfoNCE …) further improve robustness at high noise",
     "Both privileged supervision and dynamics-based pretraining are essential — validated by B1 & B2 ablations",
     "Results hold across 4 environments and 3 noise families (concat, project, nonlinear)"],
    0.6, 1.5, 12.0, 3.0,
    font_size=17, title="Takeaways")

add_rect(slide, 0.4, 4.9, 12.5, 2.2, RGBColor(0x08, 0x2C, 0x4A))
bullet_box(slide,
    ["Relax the privileged-information assumption (e.g., self-supervised or contrastive pretraining only)",
     "Extend to image-based / partial-observable environments",
     "Combine with uncertainty-aware offline RL (EDAC, IQL-uncertainty) for distribution shift robustness"],
    0.6, 5.05, 12.0, 2.0,
    font_size=16, title="Future Work", title_color=C_GOLD)

add_text(slide, "Thank you!", 0, 6.8, 13.33, 0.6,
         font_size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out_path = os.path.join(os.path.dirname(__file__), "presentation.pptx")
prs.save(out_path)
print(f"Saved: {out_path}")
